# Build ArchSim Final Review draft from template
# Preserves template styling, updates content, adds diagram placeholder slides.

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

TEMPLATE = "../template/Template.pptx"
OUTPUT = "ArchSim_Final_Review_Draft.pptx"

# Template colors / fonts
COLOR_TEXT = RGBColor(0x41, 0x3C, 0x36)
COLOR_MUTED = RGBColor(0x8A, 0x82, 0x76)
COLOR_CREAM = RGBColor(0xE7, 0xE5, 0xD9)
COLOR_ACCENT = RGBColor(0x45, 0xD1, 0xB3)  # match app

FONT_DISPLAY = "Kudryashev Display"
FONT_BODY = "Open Sauce"
FONT_BODY_BOLD = "Open Sauce Bold"
FONT_BODY_LIGHT = "Open Sauce Light"


def set_text(shape, text, font_name=None, size=None, bold=None, color=None):
    """Replace all text in a text frame with new text, preserving first paragraph style."""
    tf = shape.text_frame
    # Clear all paragraphs except first
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    # Clear runs in first paragraph
    first_p = tf.paragraphs[0]
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)
    # Add fresh run with given text
    run = first_p.add_run()
    run.text = text
    if font_name:
        run.font.name = font_name
    if size:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def set_multiline(shape, lines, font_name=None, size=None, bold=None, color=None):
    """Replace text with multiple paragraphs."""
    tf = shape.text_frame
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first_p = tf.paragraphs[0]
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)

    for i, line in enumerate(lines):
        p = first_p if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_name:
            run.font.name = font_name
        if size:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color


def set_notes(slide, notes):
    """Set speaker notes on a slide."""
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes


def duplicate_slide(prs, slide_idx):
    """Duplicate a slide by copying its XML. Returns new slide."""
    source = prs.slides[slide_idx]
    # Copy layout
    blank_layout = prs.slide_layouts[6]  # "Blank"
    new_slide = prs.slides.add_slide(blank_layout)
    # Copy all shapes from source
    for shape in source.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide


def add_textbox(slide, left_in, top_in, width_in, height_in, text,
                font=FONT_BODY, size=18, bold=False, color=COLOR_TEXT, align_center=False):
    """Add a new text box to a slide."""
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in),
                                  Inches(width_in), Inches(height_in))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align_center:
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_diagram_placeholder(slide, left_in, top_in, width_in, height_in, diagram_name):
    """Add a dashed rectangle placeholder for a diagram."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(left_in), Inches(top_in),
                                    Inches(width_in), Inches(height_in))
    # Style: dashed border, no fill
    shape.fill.background()
    line = shape.line
    line.color.rgb = COLOR_ACCENT
    line.width = Pt(1.5)
    # Set dash
    from pptx.oxml.ns import qn
    lnel = shape.line._get_or_add_ln()
    prstDash = etree.SubElement(lnel, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    # Add label
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[INSERT DIAGRAM]\n{diagram_name}"
    run.font.name = FONT_BODY_BOLD
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_ACCENT
    return shape


def add_footer_ubc(slide, y=9.9):
    """Add UBC footer to a slide."""
    tb = slide.shapes.add_textbox(Inches(6.62), Inches(y), Inches(6.77), Inches(0.42))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "THE UNIVERSITY OF BRITISH COLUMBIA"
    run.font.name = FONT_BODY_BOLD
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = COLOR_TEXT


# ============================================================
# MAIN
# ============================================================
prs = Presentation(TEMPLATE)
print(f"Opened template: {len(prs.slides)} slides")

# =========================================================
# SLIDE 1: Title — KEEP AS IS (already "Final Review", date 16/04/2026)
# =========================================================
set_notes(prs.slides[0], """Hi everyone. I'm Karl, and today I'll be presenting my final review for GP2.

The project is called ArchSim — a proof-of-concept multi-agent design system. I'll walk you through why this research matters, what I built, what it reveals about architectural practice, and where it's heading.""")

# =========================================================
# SLIDE 2: Overview — Update section list
# =========================================================
s = prs.slides[1]
for shape in s.shapes:
    if shape.has_text_frame and "Introduction" in shape.text_frame.text:
        set_multiline(shape, [
            "Introduction & Context",
            "Thesis Statement",
            "The ArchSim System",
            "What It Reveals",
            "Limitations & Future Work"
        ], font_name=FONT_BODY, size=25, color=COLOR_TEXT)
        break

set_notes(prs.slides[1], """The presentation covers five sections.

First — the context: why this research matters right now.
Second — the thesis statement.
Third — the ArchSim system itself, the proof-of-concept I built.
Fourth — what this reveals about architectural practice and authorship.
And finally — the limitations I've identified and future directions.""")

# =========================================================
# SLIDE 3: AI Evolution (image slide) — KEEP, add diagram placeholder note
# =========================================================
set_notes(prs.slides[2], """Since ChatGPT launched in November 2022, AI has evolved rapidly through four phases in architectural design.

Phase 1 — AI as Generator. Midjourney, text-to-image, concept visualization.
Phase 2 — AI as Assistant. GPT-4 multimodal, helping with reports, briefs, documents.
Phase 3 — AI as Agent. Autonomous agents beginning multi-step reasoning.
And now, Phase 4 — AI as Collaborator. Multi-agent systems. This is where my research sits.""")

# =========================================================
# SLIDE 4: Tim Fu quote — KEEP
# =========================================================
set_notes(prs.slides[3], """Tim Fu, founder of Studio Tim Fu, is one of the pioneers exploring AI in architectural design.

He puts it perfectly: "Architecture could be a collaborative process between human intention and machine intelligence, where control is neither fully surrendered nor entirely retained."

This idea of shared agency between human and machine is the philosophical foundation of my research. I'm not asking AI to replace the architect. I'm asking what happens when AI can represent the whole design team, grounded in real data.""")

# =========================================================
# SLIDE 5: Current AI tools — KEEP
# =========================================================
set_notes(prs.slides[4], """We already have powerful AI tools in architecture today — ChatGPT for writing, Midjourney for imagery, Gemini for multimodal tasks, Stable Diffusion for image generation.

These tools do photo-to-concept, live mood boarding, sketch-to-render, and interactive design.

But they all share a limitation: they respond to a single user with a single prompt producing a single output. None of them can simulate a design meeting. None of them can argue. And critically — none of them are grounded in real engineering data.""")

# =========================================================
# SLIDE 6: "What if AI could simulate the whole design team?" — KEEP
# =========================================================
set_notes(prs.slides[5], """So this is the question driving my research:

What if AI could simulate the whole design team?

Not just one voice — but a structural engineer who actually reads Karamba analysis, an MEP consultant who interprets real Ladybug energy data, a code consultant cross-referencing BCBC and Vancouver zoning. Each with different priorities. All negotiating.""")

# =========================================================
# SLIDE 7: Title divider — UPDATE if needed
# =========================================================
set_notes(prs.slides[6], """"Agentic Generative Design: Using LLM Agents as Design Collaborators"

This is the framing for my work. The key word here is COLLABORATORS — not tools, not replacements. Each agent brings its own expertise, grounded in real simulation data, into a transparent negotiation with the architect.""")

# =========================================================
# SLIDE 8: Thesis Statement — UPDATE
# =========================================================
s = prs.slides[7]
for shape in s.shapes:
    if shape.has_text_frame and "multi-agent generative" in shape.text_frame.text:
        set_text(shape,
            "This research introduces a data-grounded multi-agent generative design system in which AI "
            "agents representing Structural Engineer, MEP, Cost, and Code Consultant collaborate with "
            "the Architect through natural language — grounded in real engineering simulation "
            "(Karamba3D, Ladybug) and regulatory data (BCBC, COV zoning) — proposing a new workflow "
            "for human-AI architectural collaboration that makes the invisible negotiation of design "
            "visible, transparent, and auditable.",
            font_name=FONT_BODY, size=25, color=COLOR_TEXT)
        break

set_notes(prs.slides[7], """This is the thesis statement.

Note what's new: my system is DATA-GROUNDED. The agents are not just language models making up numbers. The Structural Engineer agent reads real Karamba3D analysis. The MEP agent reads real Ladybug energy data. The Code agent cross-references actual BCBC rules and Vancouver zoning.

The LLM provides the reasoning and argumentation layer. The domain tools provide the evidence. The human architect provides the final authority.

The goal: make the invisible negotiation of architectural design — which normally happens in emails, meetings, and redlines — visible, transparent, and auditable.""")

# =========================================================
# SLIDE 9: Design process evolution — KEEP, update notes
# =========================================================
set_notes(prs.slides[8], """To understand where this research sits, here's how design processes have evolved.

Traditional design — hand drawing and CAD — relies on human creativity and intuition. Collaboration happens through meetings and physical review. Slow iteration, no AI.

Parametric design — tools like Grasshopper — brought rule-based geometry. Faster iteration, but still script-driven, single-user, no role-based negotiation.

Agentic Generative Design — my proposal — uses natural language. The process is data-grounded agent negotiation. Multiple expert roles. Real engineering and regulatory data. Design evolves through transparent dialogue.""")

# =========================================================
# SLIDE 10: What is AI Agent? — KEEP
# =========================================================
set_notes(prs.slides[9], """Before showing you the system, let me clarify one foundational concept — what exactly IS an AI Agent, and how is it different from the AI tools we already use?""")

# =========================================================
# SLIDE 11: Normal AI vs Agent — KEEP
# =========================================================
set_notes(prs.slides[10], """On the left, normal AI. A user gives a prompt. The system executes a single action. It produces an output. Then it ends. Linear, stateless, reactive.

On the right, an AI Agent. It starts with a goal. It autonomously plans. It executes multi-step actions — searching, analyzing, verifying. Then it checks: has the goal been achieved? If not, it decomposes the task and loops back. If yes, it completes.

The key difference is the loop. An AI Agent doesn't just respond — it reasons, plans, and persists until the goal is met.

In my system, each agent does exactly this — reads the design context, pulls real data from Karamba or Ladybug or BCBC, reasons from its professional values, and proposes changes, round after round.""")

# =========================================================
# SLIDE 12: Multi-agent system — UPDATE content
# =========================================================
set_notes(prs.slides[11], """When we put multiple agents together in a shared environment, we get a Multi-Agent System.

In my system, the input is an unfiltered client brief in natural language. All agents read the same information, the same site, the same constraints.

The STRUCTURAL ENGINEER agent interprets Karamba3D forces and proposes span and lateral system.
The MEP agent interprets Ladybug energy and daylight data and proposes WWR and orientation.
The CODE agent cross-references BCBC and Vancouver zoning to check compliance.
The COST agent provides budget awareness — intentionally lighter weight so design isn't purely budget-driven.
The ARCHITECT — the human — has the final authority.

They negotiate through a multi-round protocol. The output is a negotiated design, a convergence log, and an auditable decision rationale.""")

# =========================================================
# SLIDE 13: Work in Progress divider — UPDATE to "The ArchSim System"
# =========================================================
s = prs.slides[12]
for shape in s.shapes:
    if shape.has_text_frame and "Work in Progress" in shape.text_frame.text:
        set_text(shape, "The ArchSim System",
                 font_name=FONT_DISPLAY, size=120, color=COLOR_TEXT)
        break

set_notes(prs.slides[12], """Now let me walk you through ArchSim — the proof-of-concept system I built.

Over the following slides I'll show you: how the workflow shifts agency; how Gemini and Grasshopper communicate; how agents are grounded in real data; how they negotiate; how parameters converge; and finally a live demo.""")

# Save intermediate
prs.save(OUTPUT)
print(f"Saved intermediate: {OUTPUT}")
print("\nNow adding new diagram slides...")

# ============================================================
# ADD NEW SLIDES FOR DIAGRAMS
# ============================================================
# Re-open to get clean state
prs = Presentation(OUTPUT)

# We'll insert new slides before the "Expected Final Deliverables" divider (originally slide 17).
# After our updates, the structure is:
# 1-13: intro+thesis+system intro (as above)
# 14-16: old ArchSim screens (empty slides, will repurpose)
# 17: Expected Final Deliverables divider
# 18: Future layout
# 19: Thank you

# Helper: create a content slide with title + diagram placeholder
def create_diagram_slide(prs, title_text, subtitle_text, diagram_name, insert_after_idx, notes):
    # Duplicate an existing blank-ish slide (slide 3 has just footer + background)
    new_slide = duplicate_slide(prs, 2)  # slide 3 (index 2) — blank with footer
    # Reorder: move new slide to right after insert_after_idx
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # New slide was added at end; move it
    xml_slides.remove(slides[-1])
    xml_slides.insert(insert_after_idx + 1, slides[-1])

    # Title in top-left (big Kudryashev)
    add_textbox(new_slide, left_in=1.5, top_in=0.6, width_in=17, height_in=1.2,
                text=title_text, font=FONT_DISPLAY, size=64, color=COLOR_TEXT)
    # Subtitle beneath
    add_textbox(new_slide, left_in=1.5, top_in=1.9, width_in=17, height_in=0.6,
                text=subtitle_text, font=FONT_BODY_LIGHT, size=22, color=COLOR_MUTED)
    # Diagram placeholder — large, centered
    add_diagram_placeholder(new_slide, left_in=2.5, top_in=2.9, width_in=15, height_in=6.5,
                            diagram_name=diagram_name)
    set_notes(new_slide, notes)
    return new_slide


# Slide order plan (inserting after current index 12 which is "The ArchSim System" divider):
# New slides:
#   14: Agency Redistribution (diagram 10)
#   15: App Communication Flow (diagram 02)
#   16: Agent Data Grounding (diagram 03) ⭐
#   17: System Architecture (diagram 05)
#   18: One Round of Negotiation (diagram 08)
#   19: Parameter Scope & Conflicts (diagram 09)
#   20: Convergence Across Rounds (diagram 11) ⭐
#   21: User Flow (diagram 06)
#   22: LIVE DEMO (placeholder)
#   23: What It Reveals (text slide)

# Rebuild starting fresh — we'll also remove the old empty slides 14, 15, 16 (originally mid-term screens)
# and the scenario slide 16, keeping only the essential structure.

# Strategy: add new slides after slide 13 (index 12), then remove old mid-term content slides (14-16).
# After removal, slide 17 (Expected Final Deliverables) becomes the next divider.

# We'll track insertions carefully.
# First: remove old slides 14, 15, 16 (indexes 13, 14, 15 — the empty ArchSim screens)
# Note: python-pptx doesn't have direct slide removal. We'll use xml manipulation.

def remove_slide(prs, slide_idx):
    """Remove a slide by index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[slide_idx].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[slide_idx])

# Remove old slides 14, 15, 16 (mid-term screens) — indexes 13, 14, 15
# Remove from highest index first
for idx in [15, 14, 13]:
    remove_slide(prs, idx)
print(f"After removal: {len(prs.slides)} slides")

# Now slide index 13 is "Expected Final Deliverables" (was slide 17)
# Slide 14 is the future layout (was slide 18)
# Slide 15 is Thank You (was slide 19)

# Insert new diagram slides after index 12 (the "The ArchSim System" divider)
# We want: 12 (divider), 13 (Agency), 14 (Communication), 15 (Grounding), 16 (Architecture),
#          17 (Negotiation), 18 (Scope), 19 (Convergence), 20 (User Flow), 21 (Demo), 22 (Reveals),
#          23 (Expected Final Deliverables — was index 13),
#          24 (future), 25 (thank you)

# Create new slides — each one gets inserted right after index 12 in order, so we go bottom-up
# Actually: insert at index 12+1, 12+2, etc. But duplicate_slide appends at end.
# We'll add all first then reorder.

new_slides_spec = [
    # (title, subtitle, diagram_name, notes)
    ("Agency Redistribution",
     "From linear consultant reviews to networked, data-grounded negotiation",
     "10 — Agency Redistribution (Traditional vs ArchSim)",
     """Here's the core shift. Left side — traditional workflow. The architect draws, sends to consultants, waits weeks for redlines, revises, sends to cost consultant, revises again. Linear, slow, sequential. And critically: the architect is the sole author, and the reasoning behind decisions is invisible.

Right side — ArchSim. All disciplines connect to shared design parameters simultaneously. Each agent reads real data, proposes changes, and argues. The architect decides. The loop happens in minutes, not weeks.

Authorship shifts: the architect is no longer the sole creator, but the informed decision-maker within a negotiation network. Their expertise is elevated — not replaced."""),

    ("App Communication Flow",
     "How reasoning (Gemini) and geometry (Grasshopper) connect through serve.py",
     "02 — App Communication Flow",
     """Quick note on the runtime architecture. The app uses GEMINI API — specifically gemini-3-flash — to power all agent reasoning.

The browser runs the Phaser office scene, the chat interface, and holds the current design parameters as state. When the architect or an agent changes parameters, they're POSTed to a local Python server (serve.py), which writes them to a file. Grasshopper watches that file, re-solves the geometry, runs Karamba and Ladybug, and writes the results back. The browser polls for results and updates.

A roundToken ensures the browser reads the correct solve — critical when agents make rapid changes.

Important clarification: MCP (Model Context Protocol) was used during DEVELOPMENT — I used Claude Code with Rhino/GH MCP to author the Grasshopper scripts. But the RUNTIME uses this file-based sync, not MCP."""),

    ("Agent Data Grounding",
     "The agents don't make up numbers — they interpret real engineering and regulatory data",
     "03 — Agent Data Grounding ⭐ KEY SLIDE",
     """This is the most important technical slide. It explains why my agents are not just chatbots.

STRUCTURAL ENGINEER — takes live Karamba3D analysis: beam utilization, deflection, lateral drift, member forces. The LLM's job is to interpret that data and argue for span or system changes.

MEP CONSULTANT — takes Ladybug output: solar radiation, daylight hours, energy load, carbon emission estimates. The LLM interprets the data and argues for envelope changes.

CODE CONSULTANT — cross-references actual BC Building Code, City of Vancouver zoning bylaws, and COV Open Data for site information. The LLM flags non-compliance with reasons.

COST CONSULTANT — rule-based estimates. Intentionally lighter weight. Informs but doesn't veto. I made this design decision deliberately — architectural quality should not be purely budget-driven.

The LLM provides the reasoning and argumentation layer. Real data is the evidence. Human architect is the final authority.

FUTURE: in the next phase, real consultants would define their own rubrics — max span/depth ratios, preferred lateral systems, compliance thresholds — and the AI agent becomes their digital representative. The consultant's expertise is encoded, their availability is augmented."""),

    ("System Architecture",
     "Three layers: Browser, Local Server, Rhino+Grasshopper",
     "05 — System Architecture (full technical detail)",
     """A closer look at the full tech stack. Three layers.

BROWSER — Vanilla JavaScript with ES modules. Runs the Phaser 3 office scene with first-person view, agent sprites and EasyStar pathfinding. The 5-step picker UI. The chat interface. And the Gemini API calls.

LOCAL SERVER — serve.py in Python. Provides the /gh-inputs POST endpoint and /copy-snapshot endpoint. Manages the file-based sync through .gh-inputs.json and params.json.

RHINO + GRASSHOPPER — the design engine. gh-massing.py generates geometry. gh-structure.py lays out the structural grid. gh-windows.py handles facade. Karamba3D runs structural analysis. Ladybug runs environmental analysis. gh-params-export.py captures snapshots and syncs results back.

The roundToken protocol ensures synchronization."""),

    ("One Round of Negotiation",
     "How agents consult in sequence, converge on a proposal, and re-solve the model",
     "08 — Multi-Agent Negotiation Round",
     """This is how a single round works.

The Orchestrator reads the brief and current parameters, assigns consultation order, and sets round goals.

Then each agent consults in sequence — Structural, MEP, Cost, Code — each reading the current state plus real data from their domain tool. Each proposes parameter changes.

The Planner synthesizes all proposals, resolves conflicts, and outputs a parameter change set.

Grasshopper re-solves the geometry with the new parameters. Updates the analysis. If not converged, loop back to a new round.

The ARCHITECT — the human — can intervene at any point. Walk to any agent's desk, ask questions, override decisions. The architect is always the final authority."""),

    ("Parameter Scope & Conflict Zones",
     "Where multiple agents influence the same parameter — that's where negotiation happens",
     "09 — Agent-Parameter Scope Map",
     """This diagram shows WHERE the negotiation actually happens. Each dot represents an agent's influence on a parameter.

Some parameters have only one agent — column_size is purely structural. Base shape is purely architectural. These get accepted directly.

But parameters with MULTIPLE agents — floors, WWR, structural span — these are TENSION ZONES. These are where the trade-off reasoning emerges.

On the right, three tension zones visualized. FLOORS: cost wants fewer for budget, architect wants more for GFA. WWR: MEP wants less for energy, architect wants more for daylight. STRUCTURAL SPAN: engineer wants smaller for safety, architect wants larger for open plan.

These are real architectural trade-offs that happen on every project. ArchSim makes them explicit and auditable."""),

    ("Convergence Across Rounds",
     "Proof that the system works — parameters narrow to consensus over multiple rounds",
     "11 — Convergence Timeline ⭐ KEY SLIDE",
     """This is proof that the system actually reaches consensus.

Here are four parameters over three rounds:

WWR starts at 0.60. MEP flags it as too high in round 1. Cost flags it still high in round 2. By round 3, all agents agree on 0.42. Converged.

STRUCTURAL SPAN starts at 12 meters. Structural wants smaller. Architect accepts 9m. All stable by round 3.

FLOORS starts at 16. Cost wants fewer. Code confirms FSR limit. Settles at 12.

ORIENTATION shifts 15 degrees based on MEP solar analysis.

What you see here: every change is traceable. Every decision has an author. The architect can audit: WHY did WWR drop? Because MEP's Ladybug analysis showed excessive solar gain, and cost confirmed HVAC savings. That's documentation that traditional workflows don't produce."""),

    ("User Flow",
     "Two modes: Picker (manual) or Automatic (agent-driven MDO)",
     "06 — User Journey",
     """Two modes in the app.

PICKER MODE — manual selection through 5 steps: floors, base shape, section type, plan type, window type. Each step's images reflect all prior picks. Good for exploration and education.

AUTOMATIC MODE — the architect inputs a brief, and the agents run multi-round negotiation automatically. This is the MDO mode — Multi-Disciplinary Optimization through argumentation.

Both modes feed into the same multi-round agent discussion. Both end with a final design output: negotiated parameters plus 3D model plus agent logs."""),

    ("Live Demo",
     "Walkthrough of ArchSim — picker mode and automatic negotiation",
     "[LIVE SCREEN RECORDING OR DEMO]",
     """At this point I'll run a live demo of the system.

Show:
1. Starting a new session with a residential brief
2. Picker mode — walking through the 5-step design selection
3. Agent chat — asking the Structural Engineer about lateral system
4. Automatic mode — agents running 3 rounds of negotiation
5. Final output — negotiated parameters, 3D model, full conversation log

Keep demo under 5 minutes. Focus on the agent dialogue and how decisions emerge."""),

    ("What ArchSim Reveals",
     "The architect doesn't draw buildings — they arbitrate between competing demands",
     "[TEXT SLIDE — no diagram, or use Agency Redistribution again]",
     """This is the architectural insight.

This project doesn't replace the architect — it REVEALS what the architect actually does. Senior architects stopped drawing years ago. What they actually do is WEIGH TRADE-OFFS between competing disciplinary demands.

ArchSim makes this trade-off reasoning visible for the first time. Every building that gets built is the result of hundreds of invisible negotiations — between structure and cost, between energy and glazing, between code and program. Normally these negotiations disappear into meetings and email threads. Nobody documents them.

ArchSim documents them. The architect can hand a client a convergence log: here is WHY the building is 12 floors instead of 16. Here is WHY the WWR is 42 percent. That's a new kind of design documentation — not drawings, but decision rationale.

This is the architectural contribution: a methodology where computational design preserves design LOGIC, not just design OUTCOMES."""),
]

# Create all new slides
new_slide_count = len(new_slides_spec)
created_slides = []
for title, subtitle, diag_name, notes in new_slides_spec:
    slide = create_diagram_slide(prs, title, subtitle, diag_name,
                                  insert_after_idx=12, notes=notes)
    created_slides.append(slide)

print(f"Created {new_slide_count} new slides")

# =========================================================
# UPDATE slide "Expected Final Deliverables" → "Conclusions & Future Work"
# =========================================================
# After additions and removals, structure should be:
# 0-12: original intro + thesis + system intro + system divider
# 13 to 13+new_slide_count-1: new diagram slides
# Next: old "Expected Final Deliverables" + future layout + thank you
deliverables_idx = 13 + new_slide_count
s = prs.slides[deliverables_idx]
for shape in s.shapes:
    if shape.has_text_frame and "Expected Final Deliverables" in shape.text_frame.text:
        set_text(shape, "Limitations &\nFuture Work",
                 font_name=FONT_DISPLAY, size=120, color=COLOR_TEXT)
        break

set_notes(prs.slides[deliverables_idx], """Finally, let me discuss the limitations I've identified and where this research is heading.""")

# =========================================================
# SLIDE: Future/limitations layout — ADD content
# =========================================================
future_idx = deliverables_idx + 1
s = prs.slides[future_idx]

# Clear any existing content beyond the footer and add our limitation bullet points
# Use add_textbox to place new content
add_textbox(s, left_in=1.5, top_in=0.8, width_in=17, height_in=1.3,
            text="Limitations", font=FONT_DISPLAY, size=64, color=COLOR_TEXT)

add_textbox(s, left_in=1.5, top_in=2.3, width_in=17, height_in=5.5,
            text=[
                "1. Conditional picker is pre-rendered (auto mode is live MDO)",
                "2. Single typology tested — only residential tower",
                "3. No learning across sessions — each starts fresh",
                "4. Convergence not formally guaranteed — round limits only",
                "5. Cost agent intentionally lightweight — a design choice",
                "6. Single-user — not yet multi-stakeholder",
            ], font=FONT_BODY, size=28, color=COLOR_TEXT)

add_textbox(s, left_in=1.5, top_in=8.1, width_in=17, height_in=0.8,
            text="Future: Consultant-Defined Rubrics",
            font=FONT_BODY_BOLD, size=28, bold=True, color=COLOR_ACCENT)

add_textbox(s, left_in=1.5, top_in=8.9, width_in=17, height_in=0.9,
            text="Real structural engineers, MEP consultants, and code reviewers define their own evaluation criteria. "
                 "The AI agent becomes their digital representative — expertise encoded, availability augmented.",
            font=FONT_BODY_LIGHT, size=20, color=COLOR_MUTED)

set_notes(prs.slides[future_idx], """Six limitations.

First, the conditional picker is pre-rendered — but the automatic mode runs live MDO with agents proposing real changes.

Second, single typology — only residential tower has been tested. Generalization to other building types is untested.

Third, no cross-session learning — agents don't accumulate project knowledge. Each starts fresh.

Fourth, convergence is not formally guaranteed. I use round limits to prevent infinite loops, but some conflicts may be genuinely irreconcilable.

Fifth, the cost agent is intentionally lightweight — this is a design choice, not a technical limitation. Architectural quality shouldn't be purely budget-driven.

Sixth, single-user — no multi-stakeholder collaboration yet.

The most important future direction: CONSULTANT-DEFINED RUBRICS. Currently I authored the agent evaluation criteria based on engineering knowledge. But the framework supports a future where a real structural engineer defines their own rubric — max span/depth ratios, preferred lateral systems, deflection thresholds. The AI agent becomes their digital representative. The consultant's expertise is encoded; only their availability is augmented.

This scales expertise without replacing experts.""")

# =========================================================
# Thank You slide — KEEP
# =========================================================
thank_idx = future_idx + 1
set_notes(prs.slides[thank_idx], """That's ArchSim. Thank you.

I look forward to your questions and feedback.""")

# Save final
prs.save(OUTPUT)
print(f"\nFinal draft saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
